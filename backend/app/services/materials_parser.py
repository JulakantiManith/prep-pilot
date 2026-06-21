"""Presentation materials parser for extracting content from PPT/PPTX/PDF.

Extracts slide structure (titles, bullet points, slide count) from uploaded
presentation files. This content is used during presentation evaluation to
assess how well the speaker followed their own material.

Supports:
- PDF files (extracts text per page as pseudo-slides)
- PPTX files (extracts slide titles and bullet content)
- PPT files (legacy format, treated as binary — uses filename hint only)
"""

import io
import logging
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)


@dataclass
class SlideContent:
    """Content extracted from a single slide/page."""

    index: int
    title: str = ""
    content: list[str] = field(default_factory=list)


@dataclass
class MaterialsContent:
    """Structured content extracted from presentation materials."""

    slide_count: int = 0
    slides: list[SlideContent] = field(default_factory=list)
    format: str = ""
    error: str = ""

    def to_summary(self, max_slides: int = 20) -> str:
        """Generate a text summary of the materials for AI prompts.

        Args:
            max_slides: Maximum number of slides to include in summary.

        Returns:
            Formatted text summary of slide structure and content.
        """
        if self.error:
            return f"[Materials parsing failed: {self.error}]"

        if not self.slides:
            return "[No content extracted from materials]"

        parts = [f"Presentation Materials ({self.slide_count} slides, {self.format}):"]
        for slide in self.slides[:max_slides]:
            title_str = slide.title if slide.title else "(No title)"
            parts.append(f"\n--- Slide {slide.index} ---")
            parts.append(f"Title: {title_str}")
            if slide.content:
                for point in slide.content[:8]:  # Limit bullet points per slide
                    parts.append(f"  • {point}")

        if self.slide_count > max_slides:
            parts.append(f"\n... ({self.slide_count - max_slides} more slides)")

        return "\n".join(parts)


class MaterialsParser:
    """Parses presentation materials (PDF/PPTX) to extract structured content."""

    def parse(self, file_data: bytes, filename: str) -> MaterialsContent:
        """Parse presentation materials and extract slide content.

        Args:
            file_data: Raw bytes of the uploaded file.
            filename: Original filename (used to determine format).

        Returns:
            MaterialsContent with extracted slides and metadata.
        """
        lower_name = filename.lower()

        try:
            if lower_name.endswith(".pdf"):
                return self._parse_pdf(file_data)
            elif lower_name.endswith(".pptx"):
                return self._parse_pptx(file_data)
            elif lower_name.endswith(".ppt"):
                return MaterialsContent(
                    format="ppt",
                    error="Legacy .ppt format is not supported. Please convert to .pptx.",
                )
            else:
                return MaterialsContent(
                    format="unknown",
                    error=f"Unsupported file format: {filename}",
                )
        except Exception as e:
            logger.error("Failed to parse materials '%s': %s", filename, str(e))
            return MaterialsContent(
                format=lower_name.rsplit(".", 1)[-1] if "." in lower_name else "unknown",
                error=f"Parsing failed: {str(e)}",
            )

    def _parse_pdf(self, file_data: bytes) -> MaterialsContent:
        """Extract content from PDF file (each page treated as a slide).

        Args:
            file_data: Raw PDF bytes.

        Returns:
            MaterialsContent with one SlideContent per page.
        """
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(file_data))
        slides: list[SlideContent] = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            lines = [line.strip() for line in text.split("\n") if line.strip()]

            if not lines:
                slides.append(SlideContent(index=i + 1))
                continue

            # First line is treated as the title
            title = lines[0] if lines else ""
            content = lines[1:] if len(lines) > 1 else []

            slides.append(SlideContent(
                index=i + 1,
                title=title,
                content=content,
            ))

        return MaterialsContent(
            slide_count=len(slides),
            slides=slides,
            format="pdf",
        )

    def _parse_pptx(self, file_data: bytes) -> MaterialsContent:
        """Extract content from PPTX file.

        Extracts slide titles from title placeholders and bullet points
        from content placeholders.

        Args:
            file_data: Raw PPTX bytes.

        Returns:
            MaterialsContent with structured slide data.
        """
        from pptx import Presentation
        from pptx.util import Inches  # noqa: F401

        prs = Presentation(io.BytesIO(file_data))
        slides: list[SlideContent] = []

        for i, slide in enumerate(prs.slides):
            title = ""
            content: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Check if this is the title placeholder
                    if shape.shape_id == 1 or (
                        hasattr(shape, "placeholder_format")
                        and shape.placeholder_format is not None
                        and shape.placeholder_format.idx == 0
                    ):
                        title = shape.text_frame.text.strip()
                    else:
                        # Extract all paragraphs as bullet points
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                content.append(text)

            # If no title found from placeholder, use first non-empty shape text
            if not title and content:
                title = content.pop(0)

            slides.append(SlideContent(
                index=i + 1,
                title=title,
                content=content,
            ))

        return MaterialsContent(
            slide_count=len(slides),
            slides=slides,
            format="pptx",
        )
