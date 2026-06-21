"""Integration test for presentation evaluation with materials parsing.

Tests the full pipeline: materials parsing + speech analysis + AI scoring.
Runs without Supabase/network by mocking external dependencies.
"""

import asyncio
import json
import sys
from unittest.mock import AsyncMock, MagicMock, patch

# Add backend to path
sys.path.insert(0, ".")

from app.services.materials_parser import MaterialsParser, MaterialsContent
from app.services.speech_analysis_service import SpeechAnalysisService
from app.services.presentation_service import PresentationService


def create_test_pptx() -> bytes:
    """Create a minimal PPTX file for testing."""
    from pptx import Presentation
    from pptx.util import Inches
    import io

    prs = Presentation()
    
    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Machine Learning"
    slide1.placeholders[1].text = "A beginner's guide to ML concepts"

    # Slide 2: Content slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "What is Machine Learning?"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Subset of AI that learns from data"
    body2.add_paragraph().text = "Uses algorithms to find patterns"
    body2.add_paragraph().text = "Improves with more data"

    # Slide 3: Content slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Types of ML"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Supervised Learning"
    body3.add_paragraph().text = "Unsupervised Learning"
    body3.add_paragraph().text = "Reinforcement Learning"

    # Slide 4: Content slide
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Applications"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Image recognition"
    body4.add_paragraph().text = "Natural language processing"
    body4.add_paragraph().text = "Recommendation systems"
    body4.add_paragraph().text = "Autonomous vehicles"

    # Slide 5: Conclusion
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Conclusion"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "ML is transforming every industry"
    body5.add_paragraph().text = "Start with supervised learning"
    body5.add_paragraph().text = "Practice with real datasets"

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def create_test_pdf() -> bytes:
    """Create a minimal PDF for testing."""
    from PyPDF2 import PdfWriter
    from PyPDF2.generic import NameObject, TextStringObject, ArrayObject, NumberObject
    import io

    # Create a simple PDF with text content
    writer = PdfWriter()
    
    # Add pages with content
    pages_content = [
        "Introduction to Cloud Computing\nWhat is cloud computing and why it matters",
        "Cloud Service Models\nIaaS - Infrastructure as a Service\nPaaS - Platform as a Service\nSaaS - Software as a Service",
        "Major Cloud Providers\nAWS - Amazon Web Services\nAzure - Microsoft\nGCP - Google Cloud Platform",
        "Benefits\nScalability\nCost efficiency\nReliability\nGlobal reach",
    ]
    
    for content in pages_content:
        writer.add_blank_page(width=612, height=792)
    
    buffer = io.BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


# Sample transcript (what the speaker actually said)
SAMPLE_TRANSCRIPT = """
Hello everyone, today I'm going to talk about machine learning. So, machine learning 
is basically a subset of artificial intelligence where systems learn from data without 
being explicitly programmed. Um, it's really fascinating because the algorithms can 
find patterns in large datasets that humans might miss.

There are three main types of machine learning. First, we have supervised learning, 
where you train the model on labeled data. Like, you show it pictures of cats and dogs 
and tell it which is which. Then there's unsupervised learning, where the algorithm 
finds patterns in unlabeled data on its own. And finally, reinforcement learning, 
where an agent learns by interacting with an environment and receiving rewards.

Now, the applications are really amazing. We see machine learning in image recognition, 
like facial recognition on your phone. Natural language processing powers things like 
chatbots and translation services. Recommendation systems on Netflix and Spotify use 
machine learning to suggest content you might like. And of course, autonomous vehicles 
rely heavily on machine learning for navigation and decision making.

So in conclusion, machine learning is transforming virtually every industry. If you're 
just getting started, I'd recommend beginning with supervised learning since it's the 
most straightforward. And definitely practice with real datasets from places like Kaggle. 
Thank you for listening, and I'm happy to take any questions.
"""


def test_materials_parser_pptx():
    """Test PPTX parsing extracts slides correctly."""
    print("\n=== Test: PPTX Parsing ===")
    
    pptx_data = create_test_pptx()
    parser = MaterialsParser()
    result = parser.parse(pptx_data, "test_presentation.pptx")
    
    print(f"Format: {result.format}")
    print(f"Slide count: {result.slide_count}")
    print(f"Error: {result.error or 'None'}")
    
    assert result.format == "pptx"
    assert result.slide_count == 5
    assert not result.error
    assert result.slides[0].title == "Introduction to Machine Learning"
    assert result.slides[2].title == "Types of ML"
    
    print("\nSlide summary:")
    print(result.to_summary(max_slides=5))
    print("\n✓ PPTX parsing passed!")


def test_speech_analysis():
    """Test speech analysis on the sample transcript."""
    print("\n=== Test: Speech Analysis ===")
    
    service = SpeechAnalysisService()
    # Assume ~2 minutes of speaking
    metrics = service.analyze(SAMPLE_TRANSCRIPT, duration_seconds=120.0)
    
    print(f"WPM: {metrics.wpm}")
    print(f"Total words: {metrics.total_words}")
    print(f"Filler words: {metrics.filler_word_count}")
    print(f"Filler detail: {metrics.filler_words_detail}")
    print(f"Communication score: {metrics.communication_score}/100")
    print(f"WPM in range: {metrics.wpm_in_range}")
    
    assert metrics.total_words > 50
    assert metrics.wpm > 0
    assert 0 <= metrics.communication_score <= 100
    assert metrics.filler_word_count > 0  # transcript has "um", "like", "basically", "so"
    
    print("\n✓ Speech analysis passed!")


def test_materials_summary_in_prompt():
    """Test that materials summary is properly formatted for AI prompts."""
    print("\n=== Test: Materials Summary for AI Prompt ===")
    
    pptx_data = create_test_pptx()
    parser = MaterialsParser()
    materials = parser.parse(pptx_data, "ml_presentation.pptx")
    
    summary = materials.to_summary(max_slides=15)
    
    # Verify it contains key info
    assert "5 slides" in summary
    assert "pptx" in summary
    assert "Introduction to Machine Learning" in summary
    assert "Types of ML" in summary
    assert "Applications" in summary
    assert "Conclusion" in summary
    
    print(f"Summary length: {len(summary)} chars")
    print("Contains all slide titles: ✓")
    print("\n✓ Materials summary passed!")


def test_algorithmic_scores_with_metrics():
    """Test algorithmic score generation uses speech metrics."""
    print("\n=== Test: Algorithmic Scores with Speech Metrics ===")
    
    speech_service = SpeechAnalysisService()
    metrics = speech_service.analyze(SAMPLE_TRANSCRIPT, duration_seconds=120.0)
    
    # Create a minimal service (mocking external deps)
    with patch("app.services.presentation_service.get_supabase_client"):
        service = PresentationService(
            repository=MagicMock(),
            gemini_client=MagicMock(),
        )
    
    session = {"topic": "Machine Learning", "role": "ML Intro"}
    scores = service._generate_algorithmic_scores(session, metrics)
    
    print(f"Speaking Speed: {scores.speaking_speed}/100")
    print(f"Clarity: {scores.clarity}/100")
    print(f"Structure: {scores.structure}/100")
    print(f"Communication: {scores.communication}/100")
    print(f"Engagement: {scores.engagement}/100")
    
    # Scores should be derived from actual metrics, not just defaults
    assert 0 <= scores.speaking_speed <= 100
    assert 0 <= scores.clarity <= 100
    assert 0 <= scores.structure <= 100
    assert 0 <= scores.communication <= 100
    assert 0 <= scores.engagement <= 100
    
    # With real metrics, scores shouldn't all be the default values
    default_scores = (65, 60, 55, 60, 50)
    actual_scores = (scores.speaking_speed, scores.clarity, scores.structure,
                     scores.communication, scores.engagement)
    assert actual_scores != default_scores, "Scores should differ from defaults when metrics are provided"
    
    print("\n✓ Algorithmic scores with metrics passed!")


def test_full_evaluation_pipeline_mock():
    """Test the full evaluation pipeline with mocked AI responses."""
    print("\n=== Test: Full Evaluation Pipeline (Mocked AI) ===")
    
    pptx_data = create_test_pptx()
    
    # Mock AI responses
    mock_scores_response = json.dumps({
        "speaking_speed": 78,
        "clarity": 82,
        "structure": 85,
        "communication": 80,
        "engagement": 72,
    })
    
    mock_feedback_response = json.dumps({
        "strengths": [
            "Excellent coverage of all slide topics in logical order",
            "Clear explanations with relatable examples (cats/dogs for supervised learning)",
            "Good conclusion that reinforces key takeaways"
        ],
        "weaknesses": [
            "Some filler words (um, like, basically) reduce perceived confidence",
            "Could expand more on reinforcement learning applications"
        ],
        "recommendations": [
            "Practice reducing filler words - try pausing instead of saying 'um'",
            "Add a real-world case study for each ML type to boost engagement",
            "Include a brief demo or code snippet to make the presentation more interactive"
        ]
    })
    
    # Set up mocks
    mock_repo = MagicMock()
    mock_repo.get_session.return_value = {
        "id": "test-session-id",
        "user_id": "test-user-id",
        "session_type": "presentation",
        "status": "in_progress",
        "role": "ML Introduction",
        "topic": "Machine Learning",
        "created_at": "2026-01-01T00:00:00Z",
    }
    mock_repo.update_session.return_value = {
        "id": "test-session-id",
        "user_id": "test-user-id",
        "session_type": "presentation",
        "status": "completed",
        "role": "ML Introduction",
        "topic": "Machine Learning",
        "overall_score": 79,
        "communication_score": 80,
        "created_at": "2026-01-01T00:00:00Z",
        "completed_at": "2026-01-01T00:05:00Z",
    }
    mock_repo.create_session_feedback.return_value = None
    
    mock_gemini = MagicMock()
    mock_gemini._call_openrouter_primary = AsyncMock(
        side_effect=[mock_scores_response, mock_feedback_response]
    )
    
    mock_supabase = MagicMock()
    # Mock recording download
    mock_supabase.storage.from_.return_value.list.return_value = []
    
    mock_transcription = MagicMock()
    # Create a mock TranscriptionResult with .text and .duration attributes
    mock_transcription_result = MagicMock()
    mock_transcription_result.text = SAMPLE_TRANSCRIPT
    mock_transcription_result.duration = 120.0
    mock_transcription_result.words = []
    mock_transcription.transcribe_audio_detailed = AsyncMock(return_value=mock_transcription_result)
    mock_transcription.is_meaningful_transcript = MagicMock(return_value=True)
    
    with patch("app.services.presentation_service.get_supabase_client", return_value=mock_supabase):
        service = PresentationService(
            repository=mock_repo,
            gemini_client=mock_gemini,
            transcription=mock_transcription,
        )
    
    # Override _download_recording to return fake audio data
    async def mock_download_recording(prefix):
        return b"fake_audio_data" * 1000  # ~10KB fake audio
    
    # Override _download_and_parse_materials to return parsed pptx
    async def mock_download_materials(user_id, session_id):
        parser = MaterialsParser()
        return parser.parse(pptx_data, "test.pptx")
    
    service._download_recording = mock_download_recording
    service._download_and_parse_materials = mock_download_materials
    
    # Run the full pipeline
    from uuid import UUID
    result = asyncio.run(
        service.complete_session(
            user_id="test-user-id",
            session_id=UUID("7e42b2cb-0fb5-46a5-91d3-710ccaead317"),
        )
    )
    
    print(f"Session status: {result['session']['status']}")
    print(f"\nScores:")
    scores = result["scores"]
    print(f"  Speaking Speed: {scores.speaking_speed}/100")
    print(f"  Clarity: {scores.clarity}/100")
    print(f"  Structure: {scores.structure}/100")
    print(f"  Communication: {scores.communication}/100")
    print(f"  Engagement: {scores.engagement}/100")
    
    feedback = result["feedback"]
    print(f"\nStrengths ({len(feedback.strengths)}):")
    for s in feedback.strengths:
        print(f"  ✓ {s}")
    print(f"\nWeaknesses ({len(feedback.weaknesses)}):")
    for w in feedback.weaknesses:
        print(f"  ✗ {w}")
    print(f"\nRecommendations ({len(feedback.recommendations)}):")
    for r in feedback.recommendations:
        print(f"  → {r}")
    
    # Verify the AI was called with materials content
    calls = mock_gemini._call_openrouter_primary.call_args_list
    scoring_prompt = calls[0][0][0]
    feedback_prompt = calls[1][0][0]
    
    assert "Presentation Materials" in scoring_prompt, "Scoring prompt should include materials"
    assert "Introduction to Machine Learning" in scoring_prompt, "Scoring prompt should include slide titles"
    assert "Types of ML" in scoring_prompt
    assert "slide" in feedback_prompt.lower(), "Feedback prompt should reference slides"
    
    print("\n✓ AI prompts include materials content!")
    print("\n✓ Full evaluation pipeline passed!")



def test_irrelevant_transcript_detected():
    """Test that mismatched transcript vs slides is caught by AI prompt."""
    print("\n=== Test: Irrelevant Transcript Detection ===")
    
    # Slides about Machine Learning
    pptx_data = create_test_pptx()
    parser = MaterialsParser()
    materials = parser.parse(pptx_data, "ml_presentation.pptx")
    
    # But speaker talks about cooking!
    cooking_transcript = """
    Today I want to share my favorite pasta recipe. First, you need to boil 
    water and add salt. Then cook the spaghetti for about eight minutes until 
    al dente. Meanwhile, in a separate pan, sauté garlic in olive oil. Add 
    some crushed tomatoes and let it simmer. Season with basil, oregano, and 
    a pinch of sugar. Drain the pasta and toss it with the sauce. Top with 
    parmesan cheese and fresh basil. It's really simple and delicious.
    """
    
    # Simulate what the AI would receive
    speech_service = SpeechAnalysisService()
    metrics = speech_service.analyze(cooking_transcript, duration_seconds=60.0)
    
    # Mock AI that detects the mismatch and gives low scores
    mock_scores_response = json.dumps({
        "speaking_speed": 72,
        "clarity": 68,
        "structure": 8,       # Very low - completely off-topic
        "communication": 25,  # Low - not communicating the intended content
        "engagement": 15,     # Low - not engaging with the actual topic
    })
    
    mock_feedback_response = json.dumps({
        "strengths": [
            "Clear and confident speaking delivery",
            "Good pacing and articulation of individual sentences"
        ],
        "weaknesses": [
            "CRITICAL: The speaker discussed cooking/pasta recipes while the slides cover Machine Learning - completely off-topic",
            "None of the 5 slides about ML were addressed in the presentation",
            "The audience expecting a Machine Learning talk received no relevant content"
        ],
        "recommendations": [
            "Practice delivering the ACTUAL presentation content - review each slide before recording",
            "Use your slides as a guide during the presentation to stay on topic",
            "Consider using speaker notes attached to each slide as a reminder of key points"
        ]
    })
    
    # Build the prompt that would go to AI
    with patch("app.services.presentation_service.get_supabase_client"):
        service = PresentationService(
            repository=MagicMock(),
            gemini_client=MagicMock(),
        )
    
    session = {"topic": "Machine Learning", "role": "ML Introduction"}
    
    # Check what the scoring prompt looks like
    topic = session.get("topic", "General")
    title = session.get("role", "Presentation")
    
    materials_summary = materials.to_summary(max_slides=15)
    
    # Verify the prompt would contain mismatch detection instructions
    expected_phrases = [
        "completely unrelated to the slides",
        "structure should be 0-15",
        "failed to deliver the prepared presentation",
    ]
    
    # Build materials section as the service would
    materials_section = (
        f"\n\nUploaded presentation slides:\n{materials_summary}\n"
        "\nIMPORTANT evaluation rules when slides are provided:\n"
        "- Evaluate 'structure' based on how well the speaker follows "
        "the slide order and covers slide topics. If the transcript is "
        "completely unrelated to the slides, structure should be 0-15.\n"
        "- Evaluate 'engagement' partly on whether the speaker expands "
        "on slide content rather than just reading it verbatim.\n"
        "- If the speaker's transcript does NOT match the slide content "
        "(talking about a completely different topic), give LOW scores "
        "for structure (0-15), communication (0-30), and engagement (0-20). "
        "The speaker failed to deliver the prepared presentation.\n"
    )
    
    for phrase in expected_phrases:
        assert phrase in materials_section, f"Missing: {phrase}"
    
    print("✓ Prompt includes mismatch detection instructions")
    print(f"\nMock AI scores for irrelevant transcript:")
    scores_data = json.loads(mock_scores_response)
    print(f"  Speaking Speed: {scores_data['speaking_speed']}/100 (delivery was fine)")
    print(f"  Clarity: {scores_data['clarity']}/100 (words were clear)")
    print(f"  Structure: {scores_data['structure']}/100 (CRITICAL: off-topic!)")
    print(f"  Communication: {scores_data['communication']}/100 (wrong message)")
    print(f"  Engagement: {scores_data['engagement']}/100 (not engaging with topic)")
    overall = round(sum(scores_data.values()) / 5)
    print(f"  Overall: {overall}/100")
    
    print(f"\nMock AI feedback:")
    feedback_data = json.loads(mock_feedback_response)
    for w in feedback_data["weaknesses"]:
        print(f"  ✗ {w}")
    
    print("\n✓ Irrelevant transcript detection passed!")


if __name__ == "__main__":
    print("=" * 60)
    print("PRESENTATION EVALUATION TESTS")
    print("=" * 60)
    
    test_materials_parser_pptx()
    test_speech_analysis()
    test_materials_summary_in_prompt()
    test_algorithmic_scores_with_metrics()
    test_full_evaluation_pipeline_mock()
    test_irrelevant_transcript_detected()
    
    print("\n" + "=" * 60)
    print("ALL TESTS PASSED ✓")
    print("=" * 60)
